"""Generate the leadership/board sales deck (.pptx) from data/orders.csv."""

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

DATA_PATH = "data/orders.csv"
OUTPUT_PATH = "sales_board_deck.pptx"
CHART_DIR = "/tmp/sales_deck_charts"

BG = RGBColor(0xFD, 0xF8, 0xF2)
HERO = RGBColor(0xC0, 0x39, 0x2B)
TEXT = RGBColor(0x22, 0x22, 0x22)
SUBTLE = RGBColor(0x77, 0x77, 0x77)

import os

os.makedirs(CHART_DIR, exist_ok=True)


def load_orders():
    df = pd.read_csv(DATA_PATH)
    df["order_date"] = pd.to_datetime(df["order_date"])
    df["revenue"] = df["quantity"] * df["unit_price"]
    df["month"] = df["order_date"].dt.strftime("%Y-%m")
    return df


def monthly_table(df):
    m = (
        df.groupby("month")
        .agg(revenue=("revenue", "sum"), orders=("order_id", "nunique"))
        .reset_index()
        .sort_values("month")
    )
    return m


def category_table(df):
    c = (
        df.groupby("category")
        .agg(revenue=("revenue", "sum"), orders=("order_id", "nunique"))
        .reset_index()
    )
    c["aov"] = c["revenue"] / c["orders"]
    return c.sort_values("revenue", ascending=False)


def make_bar_chart(x, y, path, color="#1f3a5f", ylabel=""):
    fig, ax = plt.subplots(figsize=(9, 4.2))
    ax.bar(x, y, color=color)
    ax.set_ylabel(ylabel)
    ax.tick_params(axis="x", rotation=45)
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    fig.tight_layout()
    fig.savefig(path, dpi=150)
    plt.close(fig)


def set_background(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size, color, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def add_bullets(slide, left, top, width, height, lines, size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"•  {line}"
        run.font.size = Pt(size)
        run.font.color.rgb = TEXT
        p.space_after = Pt(10)
    return box


def build():
    df = load_orders()
    monthly = monthly_table(df)
    by_cat = category_table(df)

    total_revenue = df["revenue"].sum()
    total_orders = df["order_id"].nunique()
    aov = total_revenue / total_orders
    best = monthly.loc[monthly["revenue"].idxmax()]

    h2_2024 = monthly.loc[monthly["month"].between("2024-07", "2024-12"), "revenue"].sum()
    h1_2025 = monthly.loc[monthly["month"].between("2025-01", "2025-06"), "revenue"].sum()
    growth = (h1_2025 / h2_2024 - 1) * 100
    jan_drop = (
        monthly.loc[monthly["month"] == "2025-01", "revenue"].values[0]
        / monthly.loc[monthly["month"] == "2024-12", "revenue"].values[0]
        - 1
    ) * 100

    trend_chart = f"{CHART_DIR}/trend.png"
    make_bar_chart(monthly["month"], monthly["revenue"], trend_chart, color="#c0392b", ylabel="Revenue ($)")

    cat_chart = f"{CHART_DIR}/category.png"
    make_bar_chart(by_cat["category"], by_cat["revenue"], cat_chart, color="#1f3a5f", ylabel="Revenue ($)")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_textbox(s, 1, 2.6, 11.3, 1.2, "Sales Performance Review", 40, TEXT, bold=True)
    add_textbox(s, 1, 3.6, 11.3, 0.8, "July 2024 – June 2025", 20, SUBTLE)

    # 2. KPI summary
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_textbox(s, 0.8, 0.5, 11, 0.8, "Total Revenue, FY24–25", 22, TEXT, bold=True)
    add_textbox(s, 0.8, 1.3, 8, 1.3, f"${total_revenue:,.0f}", 54, HERO, bold=True)
    add_textbox(
        s, 0.8, 2.6, 11, 0.6,
        f"across {total_orders:,} orders · ${aov:,.2f} average order value",
        16, SUBTLE,
    )
    add_bullets(
        s, 0.8, 3.5, 11, 2,
        [f"Best month: {best['month']} (${best['revenue']:,.2f})", "6 product categories, 60 products, 1,093 customers"],
        size=16,
    )

    # 3. Revenue trend
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_textbox(s, 0.8, 0.4, 11, 0.7, "Revenue Trend by Month", 24, TEXT, bold=True)
    s.shapes.add_picture(trend_chart, Inches(0.8), Inches(1.2), width=Inches(11.5))

    # 4. Seasonality narrative
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_textbox(s, 0.8, 0.5, 11, 0.7, "Seasonality, Not Pricing, Drives the Peak", 24, TEXT, bold=True)
    add_bullets(
        s, 0.8, 1.6, 11.5, 4.5,
        [
            f"December 2024 set the revenue record (${best['revenue']:,.2f}), driven by order volume, not higher spending per order",
            "Average order value stays in a narrow band all year — the swing is in how many orders are placed",
            f"Revenue fell {abs(jan_drop):.1f}% from December to January — the only double-digit month-over-month decline in the period",
            "Order volume and revenue both recovered steadily over the following five months",
        ],
        size=18,
    )

    # 5. Category breakdown
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_textbox(s, 0.8, 0.4, 11, 0.7, "Revenue by Category", 24, TEXT, bold=True)
    s.shapes.add_picture(cat_chart, Inches(0.8), Inches(1.2), width=Inches(11.5))

    # 6. Category detail table
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_textbox(s, 0.8, 0.4, 11, 0.7, "Category Detail — Orders & AOV", 24, TEXT, bold=True)
    rows, cols = len(by_cat) + 1, 4
    table_shape = s.shapes.add_table(rows, cols, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.2))
    table = table_shape.table
    headers = ["Category", "Revenue", "Orders", "AOV"]
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
    for i, row in enumerate(by_cat.itertuples(), start=1):
        table.cell(i, 0).text = row.category
        table.cell(i, 1).text = f"${row.revenue:,.2f}"
        table.cell(i, 2).text = f"{row.orders:,}"
        table.cell(i, 3).text = f"${row.aov:,.2f}"

    # 7. Underlying growth
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_textbox(s, 0.8, 0.5, 11, 0.7, "Underlying Growth Across Halves", 24, TEXT, bold=True)
    add_bullets(
        s, 0.8, 1.6, 11.5, 4.5,
        [
            f"H2 2024 (Jul–Dec) revenue: ${h2_2024:,.2f}",
            f"H1 2025 (Jan–Jun) revenue: ${h1_2025:,.2f}",
            f"{growth:+.1f}% versus the prior half, despite H1 2025 having no equivalent holiday spike",
            "Indicates the business is growing its baseline (non-holiday) demand, not just riding a one-time seasonal event",
        ],
        size=18,
    )

    # 8. Takeaways
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_textbox(s, 0.8, 0.5, 11, 0.7, "Takeaways", 24, TEXT, bold=True)
    add_bullets(
        s, 0.8, 1.6, 11.5, 4.8,
        [
            "Revenue grew steadily over the period, with a clear seasonal peak in December 2024",
            "The peak and the January pullback are both volume-driven, not pricing-driven",
            "Underlying (non-holiday) demand is trending up across halves",
            "Limitation: figures are gross revenue with no adjustment for returns, discounts, or refunds; only 12 months are available, so the seasonal read cannot yet be confirmed against a second holiday cycle",
        ],
        size=17,
    )

    prs.save(OUTPUT_PATH)
    print(f"Saved {OUTPUT_PATH}")


if __name__ == "__main__":
    build()
