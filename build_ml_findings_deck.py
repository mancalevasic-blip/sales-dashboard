"""Generate a leadership deck on the sales-prediction model: results, commentary,
and data-grounded suggestions to increase sales.

Reads model_metrics.json (produced by predict_sales.py) plus data/orders.csv,
and reuses the light-editorial visual style from build_sales_deck.py.
"""

import json

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

METRICS_PATH = "model_metrics.json"
DATA_PATH = "data/orders.csv"
CHART_DIR = "charts"
OUTPUT_PATH = "sales_ml_findings_deck.pptx"

BG = RGBColor(0xFD, 0xF8, 0xF2)
HERO = RGBColor(0xC0, 0x39, 0x2B)
TEXT = RGBColor(0x22, 0x22, 0x22)
SUBTLE = RGBColor(0x77, 0x77, 0x77)


def set_background(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size, color, bold=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def add_bullets(slide, left, top, width, height, lines, size=17):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"•  {line}"
        run.font.size = Pt(size)
        run.font.color.rgb = TEXT
        p.space_after = Pt(10)
    return box


def build():
    with open(METRICS_PATH) as f:
        metrics = json.load(f)

    best = next(m for m in metrics["models"] if m["model"] == metrics["best_model"])
    baseline = next(m for m in metrics["models"] if m != best)
    importances = metrics["feature_importance"]
    top_features = sorted(importances.items(), key=lambda kv: kv[1], reverse=True)[:3]

    df = pd.read_csv(DATA_PATH)
    df["order_date"] = pd.to_datetime(df["order_date"])
    df["revenue"] = df["quantity"] * df["unit_price"]
    by_cat = (
        df.groupby("category")
        .agg(revenue=("revenue", "sum"), orders=("order_id", "nunique"))
        .reset_index()
        .sort_values("revenue", ascending=False)
    )
    top_category = by_cat.iloc[0]
    bottom_category = by_cat.iloc[-1]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_textbox(s, 1, 2.6, 11.3, 1.2, "Sales Prediction Model", 40, TEXT, bold=True)
    add_textbox(s, 1, 3.6, 11.3, 0.8, "Results, Commentary & Recommendations", 20, SUBTLE)

    # 2. Model performance (hero)
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_textbox(s, 0.8, 0.5, 11, 0.8, "Model Performance", 22, TEXT, bold=True)
    add_textbox(s, 0.8, 1.3, 8, 1.3, f"{best['r2']*100:.0f}%", 54, HERO, bold=True)
    add_textbox(
        s, 0.8, 2.6, 11, 0.6,
        "of order-revenue variance explained, using only category, quantity, product, "
        "and timing — without seeing price",
        16, SUBTLE,
    )
    add_bullets(
        s, 0.8, 3.5, 11, 2.2,
        [
            f"Typical prediction error: ${best['mae']:,.2f} per order "
            f"({metrics['relative_error_pct']:.1f}% of the ${metrics['mean_revenue']:,.2f} average order)",
            f"Beats a simple baseline model (R² {baseline['r2']:.2f} → {best['r2']:.2f})",
            "Overall assessment: " + metrics["quality"].upper(),
        ],
        size=17,
    )

    # 3. Model comparison chart
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_textbox(s, 0.8, 0.4, 11, 0.7, "Model Comparison", 24, TEXT, bold=True)
    s.shapes.add_picture(f"{CHART_DIR}/model_comparison.png", Inches(1.3), Inches(1.2), width=Inches(10.5))

    # 4. What drives revenue (feature importance)
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_textbox(s, 0.8, 0.4, 11, 0.7, "What Drives Order Revenue", 24, TEXT, bold=True)
    s.shapes.add_picture(f"{CHART_DIR}/feature_importance.png", Inches(2.2), Inches(1.1), width=Inches(8.5))

    # 5. Commentary on results
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_textbox(s, 0.8, 0.5, 11, 0.7, "What This Means for the Business", 24, TEXT, bold=True)
    add_bullets(
        s, 0.8, 1.6, 11.5, 4.8,
        [
            f"Which product is ordered ({top_features[0][0]}) and how many units "
            f"({top_features[1][0] if len(top_features) > 1 else ''}) explain most of the "
            "variation in revenue — product mix and basket size matter more than category or timing",
            "Category and seasonality features contribute, but far less — revenue is driven "
            "primarily at the product level, not the category level",
            "Because the model never sees price and still predicts revenue well, product "
            "identity is acting as a reliable stand-in for price tier — a sign that pricing "
            "by product is fairly consistent and predictable today",
            f"{top_category['category']} remains the strongest category "
            f"(${top_category['revenue']:,.0f} revenue), while {bottom_category['category']} "
            f"lags furthest behind (${bottom_category['revenue']:,.0f})",
        ],
        size=17,
    )

    # 6. Recommendations
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_textbox(s, 0.8, 0.5, 11, 0.7, "Recommendations to Increase Sales", 24, TEXT, bold=True)
    add_bullets(
        s, 0.8, 1.5, 11.5, 5.2,
        [
            "Prioritize inventory and marketing spend on proven top sellers — product identity "
            "is the single biggest revenue driver, so doubling down on what already sells beats "
            "spreading effort evenly across the catalog",
            "Increase basket size — quantity is the #2 driver. Bundle offers, volume discounts, "
            "and checkout cross-sells directly target the lever with the most leverage",
            f"Grow {bottom_category['category']} toward {top_category['category']}-level "
            "performance with targeted promotions or bundling — it's the clearest underperforming "
            "category in the data",
            "Replicate the December holiday volume surge earlier in Q4 with proactive campaigns, "
            "and run a post-holiday promotion in January to soften the seasonal pullback",
            "Sustain H1 2025's underlying (non-holiday) growth with retention and loyalty programs, "
            "since that growth is organic rather than seasonal",
        ],
        size=16,
    )

    # 7. Limitations & next steps
    s = prs.slides.add_slide(blank)
    set_background(s, BG)
    add_textbox(s, 0.8, 0.5, 11, 0.7, "Limitations & Next Steps", 24, TEXT, bold=True)
    add_bullets(
        s, 0.8, 1.6, 11.5, 4.5,
        [
            "Model is trained on 12 months of data from a single year — no second holiday cycle "
            "to confirm seasonality patterns repeat",
            "Price itself was deliberately excluded as a feature, so this model explains revenue "
            "drivers, not price sensitivity or elasticity",
            "Next step: track whether the top-seller and basket-size strategies move actual "
            "revenue, then retrain the model on a second year of data to validate seasonality",
        ],
        size=18,
    )

    prs.save(OUTPUT_PATH)
    print(f"Saved {OUTPUT_PATH}")


if __name__ == "__main__":
    build()
